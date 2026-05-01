"""Extract text from Runbeck Executive Presentation .pptx."""
import json
import sys
from pathlib import Path
from pptx import Presentation

src = Path(r"C:\Users\Bogdan\Desktop\Run-Test Policies, Contracts\Policies I bookmarked - and audited, including .pptx executive presentation\Runbeck Executive Presentation (1).pptx")
out_json = Path(r"C:\Users\Bogdan\Desktop\Run-Test Policies, Contracts\validation-2026-04-27\presentation.json")
out_md = Path(r"C:\Users\Bogdan\Desktop\Run-Test Policies, Contracts\validation-2026-04-27\presentation.md")

prs = Presentation(str(src))
slides = []
for i, slide in enumerate(prs.slides, 1):
    title = ""
    body_lines = []
    notes = ""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            txt = "".join(run.text for run in para.runs).strip()
            if not txt:
                continue
            if not title and shape == slide.shapes.title:
                title = txt
            else:
                body_lines.append(txt)
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    slides.append({"slide": i, "title": title, "body": body_lines, "notes": notes})

out_json.write_text(json.dumps(slides, indent=2), encoding="utf-8")

md = ["# Runbeck Executive Presentation — Slide Text\n"]
for s in slides:
    md.append(f"\n## Slide {s['slide']}: {s['title'] or '(no title)'}\n")
    for line in s["body"]:
        md.append(f"- {line}")
    if s["notes"]:
        md.append(f"\n**Speaker Notes:**\n{s['notes']}\n")
out_md.write_text("\n".join(md), encoding="utf-8")

print(f"Extracted {len(slides)} slides")
print(f"JSON: {out_json}")
print(f"MD:   {out_md}")
