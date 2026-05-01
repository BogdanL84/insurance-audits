"""
5_Report_Builder.py — Generate email drafts, markdown reports, and slide outlines.

Three-tab layout:
  1. Email Draft — AM or client-facing
  2. Markdown Report — full detailed report
  3. Slide Outline — numbered slide-by-slide breakdown
"""

import sys
import json
import zipfile
import io
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent.parent))

import streamlit as st

st.set_page_config(
    page_title="Report Builder — Insurance Audit",
    layout="wide",
    initial_sidebar_state="expanded",
)

from datetime import date
from config import (
    CLIENTS_DIR, COLOR_NAVY, COLOR_GOOD, COLOR_BAD, COLOR_UGLY,
    BROKER_NAME, BROKER_COMPANY, BROKER_LOGO, ASSETS_DIR,
)
from core import audit_state as ast
from core.report_writer import (
    generate_email_draft,
    generate_markdown_report,
    generate_slide_outline,
    _severity_label,
)
from core.pdf_annotator import annotate_all_policies
from utils import (
    render_sidebar, require_client, render_progress_bar,
    inject_css, render_breadcrumb,
)

inject_css()
render_sidebar()

slug, client_path, state = require_client()
display_name = state.get("display_name", slug)
today_iso    = date.today().isoformat()
today_str    = date.today().strftime("%Y%m%d")

render_breadcrumb(display_name, "Report Builder")
st.title("Report Builder")
st.caption(f"**{display_name}**")
render_progress_bar(state.get("stage", "findings_reviewed"), active_step=5)
st.divider()


# ── Empty state ────────────────────────────────────────────────────
findings = state.get("findings", [])

if not findings:
    st.markdown("<br>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns([1, 2, 1])
    with c2:
        st.markdown(
            "<div style='text-align:center;padding:2rem 0'>"
            "<div style='font-size:3rem'>&#128196;</div>"
            "<h3>No findings to report</h3>"
            "<p style='color:#666'>Run the full audit analysis before building a report.</p>"
            "</div>",
            unsafe_allow_html=True,
        )
        if st.button("Run Analysis", type="primary", use_container_width=True):
            st.switch_page("pages/2_Document_Intake.py")
    st.stop()


# ── Findings summary ───────────────────────────────────────────────
ugly_ct = sum(1 for f in findings if f.get("category") == "Ugly")
bad_ct  = sum(1 for f in findings if f.get("category") == "Bad")
good_ct = sum(1 for f in findings if f.get("category") == "Good")

st.markdown(
    f"<div style='font-size:0.875rem;color:#555;margin-bottom:1rem'>"
    f"Generating reports for <strong>{len(findings)} findings</strong>: "
    f"<span style='color:{COLOR_UGLY}'>{ugly_ct} critical</span> &middot; "
    f"<span style='color:{COLOR_BAD}'>{bad_ct} bad</span> &middot; "
    f"<span style='color:{COLOR_GOOD}'>{good_ct} good</span>"
    f"</div>",
    unsafe_allow_html=True,
)

st.markdown("<br>", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════
#  PPTX BUILDER
# ══════════════════════════════════════════════════════════════════
def _build_pptx(state: dict, display_name: str, today_iso: str,
                client_path: Path, slug: str) -> tuple:
    """Build a dark-theme .pptx deck. Returns (bytes, list_of_slide_titles)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io as _io2

    # ── Palette ───────────────────────────────────────────────────
    _BG    = "1a1a2e"
    _WHITE = "FFFFFF"
    _DIMW  = "cccccc"
    _GRAY  = "888888"
    _UGLY  = "e94560"
    _BAD   = "f5a623"
    _GOOD  = "2ecc71"
    _DARK  = "0d0d1a"
    _ROW_A = "1e1e3a"
    _ROW_B = "222240"
    _HDR   = "2a2a4e"
    _SEP   = "333355"

    def _rgb(h):
        h = h.lstrip('#')
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    _blank = prs.slide_layouts[6]  # blank

    def _new_slide(bg=_BG):
        sl = prs.slides.add_slide(_blank)
        fill = sl.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(bg)
        return sl

    def _txt(sl, text, l, t, w, h, size=16, bold=False, italic=False,
             color=_WHITE, align=PP_ALIGN.LEFT, wrap=True):
        tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text or "")
        run.font.size     = Pt(size)
        run.font.bold     = bold
        run.font.italic   = italic
        run.font.color.rgb = _rgb(color)
        run.font.name     = "Calibri"
        return tf

    def _bullets(sl, lines, l, t, w, h, size=14, color=_WHITE):
        tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = str(line or "")
            run.font.size      = Pt(size)
            run.font.color.rgb = _rgb(color)
            run.font.name      = "Calibri"
        return tf

    def _rect(sl, l, t, w, h, fill, border=None):
        shp = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(fill)
        shp.line.color.rgb = _rgb(border if border else fill)
        return shp

    def _footer(sl, policy_name="", page_ref=""):
        parts = []
        if policy_name:
            parts.append(str(policy_name))
        if page_ref:
            parts.append(f"p. {page_ref}")
        if parts:
            _txt(sl, "  |  ".join(parts), 0.5, 7.1, 12.33, 0.32,
                 size=9, color=_GRAY, align=PP_ALIGN.RIGHT)

    def _divider_slide(label, subtitle, accent, bg):
        sl = _new_slide(bg)
        _rect(sl, 0, 0, 0.5, 7.5, accent)
        _txt(sl, label, 1.0, 1.7, 11.33, 2.2,
             size=80, bold=True, color=accent, align=PP_ALIGN.LEFT)
        _txt(sl, subtitle, 1.0, 3.95, 11.33, 0.8, size=24, color=_DIMW)
        return sl

    def _finding_slide(f, cat_color, cat_label):
        sl = _new_slide()
        title        = str(f.get("title", "Untitled Finding"))
        policy_quote = str(f.get("policy_quote", "") or "")[:500]
        policy_page  = str(f.get("policy_page",  "") or f.get("page_ref", "") or "")
        policy_name  = str(f.get("policy_name",  "") or "")
        plain        = str(f.get("plain_english", "") or "")[:450]
        rec          = str(f.get("recommendation", "") or "")[:220]
        sev          = f.get("severity", "")
        lik          = f.get("likelihood", "")

        # Top accent bar
        _rect(sl, 0, 0, 13.33, 0.06, cat_color)
        # Category tag + severity badge
        _txt(sl, cat_label, 0.5, 0.12, 2.5, 0.28, size=9, bold=True, color=cat_color)
        if sev or lik:
            _txt(sl, f"Likelihood {lik}/5  |  Severity {sev}/5",
                 8.5, 0.12, 4.33, 0.28, size=9, color=_GRAY, align=PP_ALIGN.RIGHT)
        # Title
        _txt(sl, title, 0.5, 0.42, 12.33, 0.75, size=22, bold=True, color=_WHITE)
        # Rule
        _rect(sl, 0.5, 1.22, 12.33, 0.03, cat_color)
        # Column headers
        _txt(sl, "WHAT THE POLICY SAYS", 0.5, 1.32, 6.0, 0.28,
             size=9, bold=True, color=cat_color)
        _txt(sl, "WHAT THIS MEANS", 7.0, 1.32, 5.83, 0.28,
             size=9, bold=True, color=cat_color)
        # Vertical separator
        _rect(sl, 6.76, 1.3, 0.03, 5.5, _SEP)

        # Left column: policy quote
        if policy_quote:
            _rect(sl, 0.5, 1.65, 6.1, 4.8, "12122a")
            _txt(sl, f'"{policy_quote}"', 0.65, 1.78, 5.8, 4.5,
                 size=11, italic=True, color=_DIMW)
        else:
            _txt(sl, "No policy quote recorded.", 0.5, 1.78, 6.1, 0.5,
                 size=11, color=_GRAY)

        # Right column: plain English + recommendation
        right_lines = []
        if plain:
            right_lines.append(plain)
        if rec:
            right_lines.append("")
            right_lines.append(f"Recommendation: {rec}")
        if right_lines:
            _bullets(sl, right_lines, 7.0, 1.65, 5.83, 4.8, size=12, color=_DIMW)

        _footer(sl, policy_name, policy_page)
        return title

    # ── Data ──────────────────────────────────────────────────────
    findings_all = state.get("findings", [])
    ugly = [f for f in findings_all if f.get("category") == "Ugly"][:6]
    bad  = [f for f in findings_all if f.get("category") == "Bad"][:5]
    good = [f for f in findings_all if f.get("category") == "Good"]
    slide_titles = []

    # ── SLIDE 1: Title ────────────────────────────────────────────
    sl = _new_slide()
    slide_titles.append("Title Slide — Insurance Program Risk Review")
    _rect(sl, 0, 0, 13.33, 0.08, _UGLY)
    _rect(sl, 0, 6.8,  13.33, 0.7,  _DARK)
    _txt(sl, display_name, 0.7, 1.8, 11.93, 1.4,
         size=44, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, "Insurance Program Risk Review", 0.7, 3.25, 11.93, 0.8,
         size=28, color=_DIMW, align=PP_ALIGN.CENTER)
    _txt(sl, today_iso, 0.7, 4.1, 11.93, 0.5,
         size=16, color=_GRAY, align=PP_ALIGN.CENTER)
    broker_line = BROKER_NAME
    if BROKER_COMPANY:
        broker_line += f"  |  {BROKER_COMPANY}"
    _txt(sl, broker_line, 0.5, 6.88, 12.33, 0.4,
         size=12, color=_DIMW, align=PP_ALIGN.CENTER)
    # Logo (if available)
    if BROKER_LOGO:
        logo_path = ASSETS_DIR / BROKER_LOGO
        if logo_path.exists():
            sl.shapes.add_picture(str(logo_path), Inches(11.5), Inches(0.15),
                                  height=Inches(0.55))

    # ── SLIDE 2: Program Snapshot ─────────────────────────────────
    sl = _new_slide()
    slide_titles.append("Program Snapshot — Policy Summary Table")
    _txt(sl, "PROGRAM SNAPSHOT", 0.5, 0.25, 12.33, 0.35,
         size=11, bold=True, color=_UGLY)
    _txt(sl, "Current Insurance Program", 0.5, 0.62, 12.33, 0.6,
         size=28, bold=True, color=_WHITE)
    _rect(sl, 0.5, 1.27, 12.33, 0.04, _UGLY)

    # Load policy analyses
    exchange_dir = client_path / "ai-exchange"
    policy_rows  = []
    if exchange_dir.exists():
        for pa in sorted(exchange_dir.glob(f"{slug}-policy-*-analysis.json")):
            try:
                data  = json.loads(pa.read_text(encoding="utf-8", errors="replace"))
                lims  = data.get("limits", {})
                # Extract premium from limits dict
                premium = "—"
                for pk in ("total_estimated_annual_premium", "annual_premium", "premium"):
                    if pk in lims:
                        v = lims[pk]
                        premium = f"${v:,.0f}" if isinstance(v, (int, float)) else str(v)
                        break
                # Top 2 non-premium limits
                skip = {"total_estimated_annual_premium", "annual_premium", "premium"}
                lim_parts = []
                for k, v in lims.items():
                    if k in skip:
                        continue
                    label = k.replace("_", " ").title()
                    val   = f"${v:,}" if isinstance(v, (int, float)) else str(v)
                    lim_parts.append(f"{label}: {val}")
                    if len(lim_parts) >= 2:
                        break
                policy_rows.append({
                    "type":    str(data.get("policy_type", "—"))[:42],
                    "carrier": str(data.get("carrier",     "—"))[:36],
                    "limits":  "; ".join(lim_parts) if lim_parts else "—",
                    "premium": premium,
                    "expiry":  str(data.get("expiry_date", "—")),
                })
            except Exception:
                pass

    # Draw table (rectangles + text boxes for full dark-theme control)
    COL_X  = [0.5,  3.2,  6.5, 10.1, 11.72]
    COL_W  = [2.65, 3.25, 3.55, 1.57, 1.61]
    HDRS   = ["Coverage Type", "Carrier", "Limits", "Premium", "Expiration"]
    HDR_H  = 0.38
    ROW_H  = 0.46
    ty     = 1.37

    _rect(sl, 0.5, ty, 12.33, HDR_H, _HDR)
    for label, x, w in zip(HDRS, COL_X, COL_W):
        _txt(sl, label, x + 0.07, ty + 0.07, w - 0.14, HDR_H - 0.1,
             size=10, bold=True, color=_UGLY)
    ty += HDR_H

    if policy_rows:
        for ri, row in enumerate(policy_rows[:8]):
            _rect(sl, 0.5, ty, 12.33, ROW_H, _ROW_A if ri % 2 == 0 else _ROW_B)
            vals = [row["type"], row["carrier"], row["limits"], row["premium"], row["expiry"]]
            for i, (val, x, w) in enumerate(zip(vals, COL_X, COL_W)):
                _txt(sl, str(val)[:58], x + 0.07, ty + 0.08, w - 0.14, ROW_H - 0.1,
                     size=9, bold=(i == 0), color=_WHITE if i == 0 else _DIMW)
            ty += ROW_H
    else:
        _txt(sl, "Policy analyses not yet available — run document analysis first.",
             0.5, 2.1, 12.33, 0.5, size=14, color=_GRAY, align=PP_ALIGN.CENTER)

    # ── UGLY findings ─────────────────────────────────────────────
    if ugly:
        _divider_slide("THE UGLY", "Critical Gaps That Must Be Fixed Now",
                       _UGLY, "2d0a14")
        slide_titles.append("Section Divider — THE UGLY")
        for f in ugly:
            t = _finding_slide(f, _UGLY, "CRITICAL")
            slide_titles.append(f"Ugly — {t}")

    # ── BAD findings ──────────────────────────────────────────────
    if bad:
        _divider_slide("THE BAD", "Issues That Need Attention",
                       _BAD, "1f1200")
        slide_titles.append("Section Divider — THE BAD")
        for f in bad:
            t = _finding_slide(f, _BAD, "GAP")
            slide_titles.append(f"Bad — {t}")

    # ── GOOD divider ──────────────────────────────────────────────
    _divider_slide("THE GOOD", "Where Your Program Is Working",
                   _GOOD, "0a1f0d")
    slide_titles.append("Section Divider — THE GOOD")

    # Good summary slide
    sl = _new_slide()
    slide_titles.append("Good Findings — What's Working")
    _rect(sl, 0, 0, 13.33, 0.06, _GOOD)
    _txt(sl, "COMPLIANT", 0.5, 0.12, 3.0, 0.28, size=9, bold=True, color=_GOOD)
    _txt(sl, "What Your Program Gets Right", 0.5, 0.42, 12.33, 0.7,
         size=28, bold=True, color=_WHITE)
    _rect(sl, 0.5, 1.18, 12.33, 0.03, _GOOD)

    good_titles = [str(f.get("title", "")) for f in good if f.get("title")][:14]
    if not good_titles:
        good_titles = ["No compliant findings recorded."]
    mid = (len(good_titles) + 1) // 2
    _bullets(sl, ["+ " + t for t in good_titles[:mid]],
             0.5, 1.4, 6.0, 5.5, size=13, color=_DIMW)
    if len(good_titles) > mid:
        _bullets(sl, ["+ " + t for t in good_titles[mid:]],
                 6.9, 1.4, 6.0, 5.5, size=13, color=_DIMW)

    # ── Our Solution ──────────────────────────────────────────────
    sl = _new_slide()
    slide_titles.append("Our Solution — Recommended Actions")
    _rect(sl, 0, 0, 13.33, 0.06, _UGLY)
    _txt(sl, "OUR SOLUTION", 0.5, 0.12, 4.0, 0.28, size=9, bold=True, color=_UGLY)
    _txt(sl, "Recommended Actions", 0.5, 0.42, 12.33, 0.7,
         size=28, bold=True, color=_WHITE)
    _rect(sl, 0.5, 1.18, 12.33, 0.03, _UGLY)

    recs = []
    for f in (ugly + bad):
        rec = str(f.get("recommendation", "") or "").strip()
        if rec and rec not in recs:
            recs.append(rec)
        if len(recs) >= 8:
            break
    if not recs:
        recs = [
            "Review current policy program with broker.",
            "Address critical coverage gaps before renewal.",
            "Schedule policy renewal strategy session.",
        ]
    _bullets(sl, [f"{i+1}. {r}" for i, r in enumerate(recs)],
             0.5, 1.38, 12.33, 5.5, size=14, color=_DIMW)

    # ── Next Steps ────────────────────────────────────────────────
    sl = _new_slide()
    slide_titles.append("Next Steps — Action Items")
    _rect(sl, 0, 0, 13.33, 0.06, _UGLY)
    _txt(sl, "NEXT STEPS", 0.5, 0.12, 4.0, 0.28, size=9, bold=True, color=_UGLY)
    _txt(sl, "Action Items", 0.5, 0.42, 12.33, 0.7,
         size=28, bold=True, color=_WHITE)
    _rect(sl, 0.5, 1.18, 12.33, 0.03, _UGLY)

    for i, placeholder in enumerate([
        "Action Item 1: _______________________________________________",
        "Action Item 2: _______________________________________________",
        "Action Item 3: _______________________________________________",
    ]):
        ty2 = 1.55 + i * 1.5
        _rect(sl, 0.5, ty2, 12.33, 1.2, _ROW_A)
        _txt(sl, placeholder, 0.7, ty2 + 0.38, 12.0, 0.5, size=16, color=_GRAY)

    # ── Serialize ─────────────────────────────────────────────────
    buf = _io2.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue(), slide_titles


# ══════════════════════════════════════════════════════════════════
#  THREE TABS
# ══════════════════════════════════════════════════════════════════
tab_email, tab_report, tab_pptx, tab_pdfs, tab_export = st.tabs([
    "Email Draft",
    "Markdown Report",
    "PowerPoint",
    "Annotated Policies",
    "Export All",
])


# ══════════════════════════════════════════════════════════════════
#  TAB 1: EMAIL DRAFT
# ══════════════════════════════════════════════════════════════════
with tab_email:
    st.subheader("Email Draft")
    st.caption(
        "Account Manager version has full technical detail. "
        "Client / CFO version is plain-English for a business owner."
    )
    st.markdown("<br>", unsafe_allow_html=True)

    col_sel, col_gen, _ = st.columns([2, 2, 3])
    with col_sel:
        recipient = st.radio(
            "Recipient",
            options=["Account Manager", "Client / CFO"],
            horizontal=True,
            key="email_recipient_radio",
        )
        rec_type = "am" if recipient == "Account Manager" else "client"

    email_cache_key = f"email_draft_{slug}_{rec_type}"
    with col_gen:
        st.markdown("<br>", unsafe_allow_html=True)
        regen = st.button("Generate Email", key="regen_email_btn", type="primary",
                          use_container_width=True)

    if regen or email_cache_key not in st.session_state:
        st.session_state[email_cache_key] = generate_email_draft(state, rec_type)

    email_text = st.session_state.get(email_cache_key, "")
    if email_text:
        st.text_area(
            label  = "Email Draft — copy and paste into your email client",
            value  = email_text,
            height = 400,
            key    = f"email_display_{rec_type}",
        )
        st.download_button(
            label     = "Download .txt",
            data      = email_text.encode("utf-8"),
            file_name = f"{slug}-email-{rec_type}-{today_str}.txt",
            mime      = "text/plain",
            key       = f"dl_email_{rec_type}",
        )


# ══════════════════════════════════════════════════════════════════
#  TAB 2: MARKDOWN REPORT
# ══════════════════════════════════════════════════════════════════
with tab_report:
    st.subheader("Markdown Report")
    st.caption(
        "Full detailed audit report with all findings, contract quotes, "
        "policy quotes, gap analysis, and recommendations."
    )
    st.markdown("<br>", unsafe_allow_html=True)

    report_cache_key = f"md_report_{slug}"
    col_gen_r, _ = st.columns([2, 5])
    with col_gen_r:
        regen_r = st.button("Generate Report", key="regen_report_btn", type="primary",
                            use_container_width=True)

    if regen_r or report_cache_key not in st.session_state:
        st.session_state[report_cache_key] = generate_markdown_report(state)

    report_text = st.session_state.get(report_cache_key, "")
    if report_text:
        st.caption(
            f"{len(report_text):,} characters &middot; {len(report_text.split()):,} words",
            unsafe_allow_html=True,
        )
        st.text_area(label="Markdown Report", value=report_text, height=450, key="report_display")

        report_filename = f"audit-report-{today_str}.md"
        dl_col, save_col, _ = st.columns([2, 2, 3])
        with dl_col:
            st.download_button(
                label     = "Download .md",
                data      = report_text.encode("utf-8"),
                file_name = report_filename,
                mime      = "text/markdown",
                key       = "dl_report",
                use_container_width=True,
            )
        with save_col:
            if st.button("Save to Output Folder", key="save_report_btn", use_container_width=True):
                output_dir = client_path / "output"
                output_dir.mkdir(parents=True, exist_ok=True)
                (output_dir / report_filename).write_text(report_text, encoding="utf-8")
                st.success(f"Saved to `output/{report_filename}`")


# ══════════════════════════════════════════════════════════════════
#  TAB 3: POWERPOINT
# ══════════════════════════════════════════════════════════════════
with tab_pptx:
    st.subheader("PowerPoint Presentation")
    st.caption(
        "Dark-theme .pptx deck — The Ugly / The Bad / The Good structure, "
        "two-column finding slides with policy quotes, and a program snapshot table."
    )
    st.markdown("<br>", unsafe_allow_html=True)

    _pptx_key   = f"pptx_deck_{slug}"
    _titles_key = f"pptx_titles_{slug}"

    col_gen_p, _ = st.columns([2, 5])
    with col_gen_p:
        _gen_pptx = st.button("Generate PowerPoint", key="gen_pptx_btn",
                               type="primary", use_container_width=True)

    if _gen_pptx:
        try:
            with st.spinner("Building deck..."):
                _pptx_bytes, _slide_titles = _build_pptx(
                    state, display_name, today_iso, client_path, slug
                )
            _output_dir_p = client_path / "output"
            _output_dir_p.mkdir(parents=True, exist_ok=True)
            _pptx_fname = f"{slug}-risk-review-{today_str}.pptx"
            (_output_dir_p / _pptx_fname).write_bytes(_pptx_bytes)
            st.session_state[_pptx_key]   = _pptx_bytes
            st.session_state[_titles_key] = _slide_titles
            st.success(f"Saved to `output/{_pptx_fname}` ({len(_pptx_bytes)//1024} KB)")
        except ImportError:
            st.error(
                "python-pptx is not installed. "
                "Run `pip install python-pptx` and restart the app."
            )
        except Exception as _exc:
            st.error(f"Failed to build presentation: {_exc}")

    _pptx_bytes   = st.session_state.get(_pptx_key)
    _slide_titles = st.session_state.get(_titles_key, [])

    if _pptx_bytes:
        _pptx_fname = f"{slug}-risk-review-{today_str}.pptx"
        st.download_button(
            label     = "Download .pptx",
            data      = _pptx_bytes,
            file_name = _pptx_fname,
            mime      = "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key       = "dl_pptx_btn",
        )

        if _slide_titles:
            _PREVIEW_COLORS = {
                "Title":   "#1A237E",
                "Program": "#37474F",
                "Section": "#555577",
                "Ugly":    "#e94560",
                "Bad":     "#f5a623",
                "Good":    "#2ecc71",
                "Our":     "#1A237E",
                "Next":    "#1A237E",
            }
            st.markdown("<br>", unsafe_allow_html=True)
            st.markdown(f"**Slide preview** — {len(_slide_titles)} slides")
            for _i, _t in enumerate(_slide_titles, 1):
                _first = _t.split()[0] if _t.split() else ""
                _clr   = _PREVIEW_COLORS.get(_first, "#37474F")
                st.markdown(
                    f"<div style='display:flex;align-items:center;margin:3px 0'>"
                    f"<span style='background:{_clr};color:#fff;border-radius:3px;"
                    f"padding:1px 8px;font-size:0.78rem;font-weight:700;min-width:30px;"
                    f"text-align:center;margin-right:10px'>{_i}</span>"
                    f"<span style='font-size:0.875rem;color:#333'>{_t}</span>"
                    f"</div>",
                    unsafe_allow_html=True,
                )


# ══════════════════════════════════════════════════════════════════
#  TAB 4: ANNOTATED POLICIES
# ══════════════════════════════════════════════════════════════════
with tab_pdfs:
    st.subheader("Annotated Policy PDFs")
    st.caption(
        "Each policy PDF marked up with highlights (yellow/orange/red by severity), "
        "sticky-note annotations, bookmarks, and a cover page. "
        "The original PDF is never modified."
    )
    st.markdown("<br>", unsafe_allow_html=True)

    policies_dir = client_path / "policies"
    output_dir   = client_path / "output"
    exchange_dir = client_path / "ai-exchange"

    def _load_policy_analyses_rb(exchange_dir: Path, slug: str) -> list:
        analyses = []
        for p in exchange_dir.glob(f"{slug}-policy-*-analysis.json"):
            try:
                analyses.append(json.loads(p.read_text(encoding="utf-8", errors="replace")))
            except Exception:
                pass
        return analyses

    # ── Existing annotated PDFs ────────────────────────────────────
    existing = sorted(output_dir.glob("*-AUDITED.pdf")) if output_dir.exists() else []

    if existing:
        st.markdown("**Available annotated PDFs:**")
        st.markdown("<br>", unsafe_allow_html=True)
        for ep in existing:
            dl_col, info_col = st.columns([2, 3])
            with dl_col:
                st.download_button(
                    label             = f"Download {ep.name}",
                    data              = ep.read_bytes(),
                    file_name         = ep.name,
                    mime              = "application/pdf",
                    key               = f"rb_dl_{ep.stem}",
                    use_container_width = True,
                    type              = "primary",
                )
            with info_col:
                size_kb = ep.stat().st_size / 1024
                st.caption(
                    f"{ep.name}  ·  {size_kb:.0f} KB"
                )
        st.markdown("<br>", unsafe_allow_html=True)
        st.caption("Re-generate to update annotations after re-running analysis.")
        st.markdown("<br>", unsafe_allow_html=True)

    # ── Generate button ────────────────────────────────────────────
    has_policies = policies_dir.exists() and any(policies_dir.glob("*.pdf"))

    if not has_policies:
        st.warning("No policy PDFs found. Upload policies in Document Intake first.")
    else:
        btn_label = "Re-generate Annotated PDFs" if existing else "Generate Annotated PDFs"
        if st.button(btn_label, type="primary" if not existing else "secondary",
                     key="rb_generate_pdfs_btn"):
            policy_analyses = _load_policy_analyses_rb(exchange_dir, slug)
            prog            = st.progress(0.0, text="Annotating policies...")
            log_box         = st.container()

            results = annotate_all_policies(
                policies_dir    = policies_dir,
                findings        = findings,
                policy_analyses = policy_analyses,
                client_name     = display_name,
                output_dir      = output_dir,
            )

            completed, failed = [], []
            for i, (fname, out_path, n_f, err) in enumerate(results):
                prog.progress(
                    (i + 1) / max(len(results), 1),
                    text=f"Annotated {fname}...",
                )
                if err:
                    failed.append((fname, err))
                    with log_box:
                        st.warning(f"{fname}: {err}")
                else:
                    completed.append((fname, out_path, n_f))
                    with log_box:
                        st.success(f"{fname}: {n_f} findings annotated.")

            prog.progress(1.0, text="Done.")

            if completed:
                st.markdown("<br>", unsafe_allow_html=True)
                for fname, out_path, n_f in completed:
                    st.download_button(
                        label             = f"Download {Path(out_path).name}",
                        data              = Path(out_path).read_bytes(),
                        file_name         = Path(out_path).name,
                        mime              = "application/pdf",
                        key               = f"rb_dl_new_{Path(out_path).stem}",
                        use_container_width = True,
                    )


# ══════════════════════════════════════════════════════════════════
#  TAB 5: EXPORT ALL
# ══════════════════════════════════════════════════════════════════

def _slides_to_markdown(slides: list) -> str:
    """Convert slide list to readable markdown."""
    lines = []
    for slide in slides:
        num     = slide.get("number", "?")
        section = slide.get("section", "").replace("_", " ").title()
        title   = slide.get("title", "")
        content = slide.get("content", "")
        notes   = slide.get("notes", "")
        lines.append(f"## Slide {num:02d} — {section}: {title}")
        if content:
            lines.append(content)
        if notes:
            lines.append(f"\n*Speaker Notes: {notes}*")
        lines.append("")
    return "\n".join(lines)


def _build_one_page_summary(state: dict, display_name: str, today_iso: str) -> str:
    """Build a one-page markdown summary of the audit."""
    findings = state.get("findings", [])
    ugly = [f for f in findings if f.get("category") == "Ugly"]
    bad  = [f for f in findings if f.get("category") == "Bad"]
    good = [f for f in findings if f.get("category") == "Good"]
    info = state.get("client_info", {})

    lines = [
        f"# Audit Summary — {display_name}",
        f"**Date:** {today_iso}",
        f"**Industry:** {info.get('industry', '—')}",
        "",
        "---",
        "",
        f"## Findings Overview",
        f"| Category | Count |",
        f"|----------|-------|",
        f"| Critical (Ugly) | {len(ugly)} |",
        f"| Gaps (Bad) | {len(bad)} |",
        f"| Compliant (Good) | {len(good)} |",
        f"| **Total** | **{len(findings)}** |",
        "",
        "---",
        "",
        "## Critical Issues",
    ]
    for f in ugly[:10]:
        sev = f.get("severity", 0)
        lik = f.get("likelihood", 0)
        lines.append(
            f"### {f.get('title', 'Untitled')}"
        )
        if sev or lik:
            lines.append(f"*Likelihood: {lik}/5 · Severity: {sev}/5*")
        pe = f.get("plain_english", "")
        if pe:
            lines.append(pe)
        rec = f.get("recommendation", "")
        if rec:
            lines.append(f"**Recommendation:** {rec}")
        lines.append("")
    if not ugly:
        lines.append("*No critical issues found.*\n")

    lines += ["---", "", "## Gaps to Address"]
    for f in bad[:10]:
        lines.append(f"- **{f.get('title', '')}** — {f.get('plain_english', '')[:200]}")
    if not bad:
        lines.append("*No gaps found.*")

    lines += ["", "---", f"*Generated by Insurance Audit System*"]
    return "\n".join(lines)


with tab_export:
    st.subheader("Export Complete Audit Package")
    st.caption(
        "Download all reports, emails, annotated PDFs, and supplemental files "
        "in a single ZIP. Any outputs not yet generated are created automatically."
    )
    st.markdown("<br>", unsafe_allow_html=True)

    _exchange_dir = client_path / "ai-exchange"
    _output_dir   = client_path / "output"

    # Show what will be included
    _sp_file = _exchange_dir / "strategic-plan.json"
    _cr_file = _exchange_dir / "client-research.md"
    _pdfs    = sorted(_output_dir.glob("*-AUDITED.pdf")) if _output_dir.exists() else []

    _include_lines = [
        "- `audit-report.md` — full findings report",
        "- `email-draft-AM.txt` — account manager email",
        "- `email-draft-CFO.txt` — client-facing email",
        "- `slide-outline.md` — presentation outline",
        "- `one-page-summary.md` — executive summary",
    ]
    if _sp_file.exists():
        _include_lines.append("- `strategic-plan.md` — strategic advisory plan")
    if _cr_file.exists():
        _include_lines.append("- `client-research.md` — client web research")
    for _p in _pdfs:
        _include_lines.append(f"- `annotated-pdfs/{_p.name}`")

    with st.expander("What's included", expanded=True):
        st.markdown("\n".join(_include_lines))

    st.markdown("<br>", unsafe_allow_html=True)

    _export_key = f"export_zip_{slug}"

    _col_btn, _ = st.columns([3, 4])
    with _col_btn:
        if st.button(
            "Prepare Complete Audit Package",
            key="prepare_zip_btn",
            type="primary",
            use_container_width=True,
        ):
            with st.spinner("Building package..."):
                _buf = io.BytesIO()
                with zipfile.ZipFile(_buf, "w", zipfile.ZIP_DEFLATED) as _zf:

                    # Email AM
                    _e_am = (
                        st.session_state.get(f"email_draft_{slug}_am")
                        or generate_email_draft(state, "am")
                    )
                    _zf.writestr("email-draft-AM.txt", _e_am or "")

                    # Email Client/CFO
                    _e_cl = (
                        st.session_state.get(f"email_draft_{slug}_client")
                        or generate_email_draft(state, "client")
                    )
                    _zf.writestr("email-draft-CFO.txt", _e_cl or "")

                    # Markdown report
                    _rpt = (
                        st.session_state.get(f"md_report_{slug}")
                        or generate_markdown_report(state)
                    )
                    _zf.writestr("audit-report.md", _rpt or "")

                    # Slide outline
                    _sl = (
                        st.session_state.get(f"slide_outline_{slug}")
                        or generate_slide_outline(state)
                    )
                    if _sl:
                        _zf.writestr("slide-outline.md", _slides_to_markdown(_sl))

                    # One-page summary
                    _zf.writestr(
                        "one-page-summary.md",
                        _build_one_page_summary(state, display_name, today_iso),
                    )

                    # Strategic plan (if generated)
                    if _sp_file.exists():
                        try:
                            _sp_data = json.loads(_sp_file.read_text(encoding="utf-8"))
                            _sp_raw  = _sp_data.get("raw", "")
                            if _sp_raw:
                                _zf.writestr("strategic-plan.md", _sp_raw)
                        except Exception:
                            pass

                    # Client research (if generated)
                    if _cr_file.exists():
                        _zf.writestr(
                            "client-research.md",
                            _cr_file.read_text(encoding="utf-8"),
                        )

                    # Annotated PDFs
                    for _pdf in _pdfs:
                        try:
                            _zf.writestr(
                                f"annotated-pdfs/{_pdf.name}",
                                _pdf.read_bytes(),
                            )
                        except Exception:
                            pass

                _buf.seek(0)
                st.session_state[_export_key] = _buf.getvalue()

    # Show download button if package is ready
    _zip_data = st.session_state.get(_export_key)
    if _zip_data:
        _zip_name = f"{slug}-audit-package-{today_str}.zip"
        _size_kb  = len(_zip_data) / 1024
        st.markdown("<br>", unsafe_allow_html=True)
        st.success(f"Package ready — {_size_kb:.0f} KB")
        st.download_button(
            label             = f"Download {_zip_name}",
            data              = _zip_data,
            file_name         = _zip_name,
            mime              = "application/zip",
            key               = "dl_zip_btn",
            use_container_width = True,
        )
        if st.button("Clear & Rebuild", key="clear_zip_btn"):
            del st.session_state[_export_key]
            st.rerun()


# ══════════════════════════════════════════════════════════════════
#  MARK COMPLETE + NAV
# ══════════════════════════════════════════════════════════════════
st.markdown("<br>", unsafe_allow_html=True)
st.divider()

complete_col, nav_col = st.columns([2, 5])
with complete_col:
    if state.get("stage") != "output_generated":
        if st.button("Mark Audit Complete", key="mark_complete_btn", use_container_width=True):
            state["stage"] = "output_generated"
            ast.save(client_path, state)
            st.success("Audit marked as complete.")
            st.rerun()
    else:
        st.markdown(
            "<div style='padding:0.5rem;background:#E8F5E9;border-radius:6px;"
            "font-size:0.875rem;color:#2E7D32'>&#10003; Audit is complete</div>",
            unsafe_allow_html=True,
        )

with nav_col:
    nl, nr = st.columns(2)
    with nl:
        if st.button("Back to Findings", use_container_width=True):
            st.switch_page("pages/3_Findings_Dashboard.py")
    with nr:
        if st.button("Dashboard", use_container_width=True):
            st.switch_page("app.py")
